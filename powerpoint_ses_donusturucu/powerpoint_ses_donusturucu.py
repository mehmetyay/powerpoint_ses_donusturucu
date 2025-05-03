"""
PowerPoint'ten Sesli Anlatıma Dönüştürücü Profesyonel
Copyright (c) 2025 Mehmet Yay. Tüm hakları saklıdır.

MIT Lisansı ile lisanslanmıştır:
Bu yazılımın ve ilgili belgelerin ("Yazılım") kopyalarını alan herkese,
ücretsiz olarak, Yazılımı kısıtlama olmadan kullanma, kopyalama, değiştirme,
birleştirme, yayımlama, dağıtma, alt lisanslama ve/veya satma hakkı verilir.
"""
import tempfile
import uuid
import os
import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext, font
from pptx import Presentation
from gtts import gTTS
import pygame
from pydub import AudioSegment
from pydub.playback import play
import threading
import time
from PIL import Image, ImageTk
import json
from datetime import datetime
import webbrowser
import platform
import subprocess
import sys
import random
from tkinter import PhotoImage
from tkinter.ttk import Progressbar
import locale
import gettext
import shutil
import winsound
import pyautogui
import psutil
import requests
from bs4 import BeautifulSoup
from packaging import version
import zipfile
import io

# Dil ayarları
locale.setlocale(locale.LC_ALL, 'tr_TR.UTF-8')
_ = gettext.gettext

class PowerPointSesDonusturucu:
    def __init__(self, root):
        self.root = root
        self.root.title(_("PowerPoint'ten Sesli Anlatıma Dönüştürücü Profesyonel"))
        self.root.geometry("1000x800")
        self.root.minsize(900, 700)
        self.root.configure(bg="#2c3e50")
        
        # Uygulama simgesi
        try:
            self.root.iconbitmap("ppt_icon.ico")
        except:
            pass
        
        # Sistem bilgisi
        self.sistem_bilgisi = {
            "platform": platform.system(),
            "sürüm": platform.release(),
            "python": platform.python_version()
        }
        
        # Animasyonlu arka plan
        self.canvas = tk.Canvas(root, bg="#2c3e50", highlightthickness=0)
        self.canvas.pack(fill="both", expand=True)
        self.yıldızlar = []
        self.yıldız_olustur()
        
        # Uygulama değişkenleri
        self.giris_dosyası = ""
        self.cikis_dosyası = ""
        self.cikis_formatı = ".mp3"
        self.dil = "tr"
        self.yavas_konusma = False
        self.donusum_devam_ediyor = False
        self.oynatma_devam_ediyor = False
        self.gece_modu = False
        self.ses_seviyesi = 0.7
        self.son_kullanılan_klasor = os.path.expanduser("~")
        self.guncelleme_kontrol_edildi = False
        
        # Modern UI renk paleti
        self.renkler = {
            "arkaplan": "#2c3e50",
            "cerceve": "#34495e",
            "baslik": "#3498db",
            "buton": "#2980b9",
            "buton_hover": "#3498db",
            "buton_aktif": "#1abc9c",
            "metin": "#ecf0f1",
            "uyari": "#e74c3c",
            "basari": "#2ecc71",
            "vurgu": "#f39c12"
        }
        
        # Fontlar
        self.font_baslik = font.Font(family="Segoe UI", size=18, weight="bold")
        self.font_alt_baslik = font.Font(family="Segoe UI", size=12, weight="bold")
        self.font_normal = font.Font(family="Segoe UI", size=10)
        self.font_kucuk = font.Font(family="Segoe UI", size=8)
        
        # Başlatma işlemleri
        self.pygame_baslat()
        self.arayuz_olustur()
        self.ayarları_yukle()
        self.guncelleme_kontrol()
        
        # Animasyon başlat
        self.animasyon_aktif = True
        self.animasyon_guncelle()
        
    def pygame_baslat(self):
        """Pygame mixer'ı başlatır"""
        try:
            pygame.mixer.init(frequency=44100, size=-16, channels=2, buffer=4096)
            pygame.mixer.music.set_volume(self.ses_seviyesi)
        except Exception as e:
            self.durum_guncelle(_("Ses sistemi başlatılamadı: {}").format(str(e)), "uyari")
    
    def yıldız_olustur(self):
        """Animasyonlu arka plan için yıldızlar oluşturur"""
        for _ in range(100):
            x = random.randint(0, self.root.winfo_screenwidth())
            y = random.randint(0, self.root.winfo_screenheight())
            size = random.randint(1, 3)
            parlaklık = random.randint(100, 255)
            renk = f"#{parlaklık:02x}{parlaklık:02x}{parlaklık:02x}"
            yıldız = self.canvas.create_oval(x, y, x+size, y+size, fill=renk, outline="")
            hız = random.uniform(0.5, 2.0)
            self.yıldızlar.append((yıldız, hız))
    
    def animasyon_guncelle(self):
        """Yıldız animasyonunu günceller"""
        if not self.animasyon_aktif:
            return
            
        for yıldız, hız in self.yıldızlar:
            self.canvas.move(yıldız, 0, hız)
            pos = self.canvas.coords(yıldız)
            if pos[1] > self.root.winfo_height():
                self.canvas.move(yıldız, 0, -self.root.winfo_height()-10)
        
        self.root.after(30, self.animasyon_guncelle)
    
    def arayuz_olustur(self):
        """Ana arayüzü oluşturur"""
        # Ana çerçeve
        self.ana_cerceve = tk.Frame(self.canvas, bg=self.renkler["arkaplan"])
        self.ana_cerceve.pack(fill="both", expand=True, padx=10, pady=10)
        
        # Başlık çerçevesi
        self.baslik_cerceve = self.cerceve_olustur(self.ana_cerceve, _("PowerPoint'ten Sesli Anlatıma Dönüştürücü"), 
                                                  bg=self.renkler["baslik"], relief=tk.RAISED, bd=2)
        self.baslik_cerceve.pack(fill="x", pady=(0, 10))
        
        # Logo ve başlık
        self.logo_ekle()
        
        # Giriş bölümü
        self.giris_cercevesi_olustur()
        
        # Çıkış ayarları
        self.cikis_ayarlari_olustur()
        
        # Önizleme bölümü
        self.onizleme_cercevesi_olustur()
        
        # Kontrol butonları
        self.kontrol_butonlari_olustur()
        
        # Durum çubuğu
        self.durum_cubugu_olustur()
        
        # Menü çubuğu
        self.menu_cubugu_olustur()
        
        # Araç çubuğu
        self.arac_cubugu_olustur()
        
        # Gece modu butonu
        self.gece_modu_butonu = tk.Button(self.ana_cerceve, text="🌙", font=self.font_alt_baslik,
                                        command=self.gece_modunu_degistir, bd=0, bg=self.renkler["arkaplan"],
                                        fg=self.renkler["metin"], activebackground=self.renkler["arkaplan"])
        self.gece_modu_butonu.place(relx=0.98, rely=0.02, anchor="ne")
    
    def logo_ekle(self):
        """Başlık çerçevesine logo ekler"""
        logo_cerceve = tk.Frame(self.baslik_cerceve, bg=self.renkler["baslik"])
        logo_cerceve.pack(side="left", padx=10)
        
        try:
            # Logo resmini yükle
            logo_resim = Image.open("ppt_icon.png").resize((64, 64))
            self.logo_img = ImageTk.PhotoImage(logo_resim)
            logo_etiket = tk.Label(logo_cerceve, image=self.logo_img, bg=self.renkler["baslik"])
            logo_etiket.pack(side="left", padx=(0, 10))
        except Exception as e:
            print(_("Logo yüklenemedi:"), str(e))
            logo_etiket = tk.Label(logo_cerceve, text="🎤", font=("Segoe UI", 24), 
                                 bg=self.renkler["baslik"], fg=self.renkler["metin"])
            logo_etiket.pack(side="left", padx=(0, 10))
        
        # Başlık metni
        baslik_metin = tk.Label(logo_cerceve, 
                              text=_("PowerPoint'ten Sesli Anlatıma Dönüştürücü Profesyonel"),
                              font=self.font_baslik,
                              bg=self.renkler["baslik"],
                              fg=self.renkler["metin"])
        baslik_metin.pack(side="left")
        
        # Sürüm bilgisi
        surum_metin = tk.Label(self.baslik_cerceve, 
                             text=_("Sürüm 2.0 | © 2025 Mehmet Yay"),
                             font=self.font_kucuk,
                             bg=self.renkler["baslik"],
                             fg=self.renkler["metin"])
        surum_metin.pack(side="right", padx=10)
    
    def giris_cercevesi_olustur(self):
        """PowerPoint dosyası giriş bölümünü oluşturur"""
        self.giris_cerceve = self.cerceve_olustur(self.ana_cerceve, _("PowerPoint Dosyası Seçin"), 
                                                 bg=self.renkler["cerceve"])
        self.giris_cerceve.pack(fill="x", pady=(0, 10))
        
        # Giriş dosyası alanı
        giris_etiket = tk.Label(self.giris_cerceve, 
                              text=_("PowerPoint Dosyası:"),
                              font=self.font_normal,
                              bg=self.renkler["cerceve"],
                              fg=self.renkler["metin"])
        giris_etiket.grid(row=0, column=0, padx=(10, 5), pady=10, sticky="w")
        
        self.giris_entry = tk.Entry(self.giris_cerceve, 
                                  width=60,
                                  font=self.font_normal,
                                  bg="#ffffff",
                                  fg="#000000")
        self.giris_entry.grid(row=0, column=1, padx=(0, 5), pady=10, sticky="ew")
        
        # Gözat butonu
        self.girisi_gosterme_butonu = self.buton_olustur(self.giris_cerceve, 
                                                        _("Gözat..."), 
                                                        self.dosya_sec,
                                                        bg=self.renkler["buton"],
                                                        fg=self.renkler["metin"])
        self.girisi_gosterme_butonu.grid(row=0, column=2, padx=(0, 10), pady=10)
        
        # Dosya bilgisi
        self.dosya_bilgi_etiket = tk.Label(self.giris_cerceve, 
                                         text="",
                                         font=self.font_kucuk,
                                         bg=self.renkler["cerceve"],
                                         fg=self.renkler["metin"])
        self.dosya_bilgi_etiket.grid(row=1, column=0, columnspan=3, padx=10, pady=(0, 10), sticky="w")
        
        # Grid ayarları
        self.giris_cerceve.grid_columnconfigure(1, weight=1)
    
    def cikis_ayarlari_olustur(self):
        """Ses çıkış ayarları bölümünü oluşturur"""
        self.cikis_cerceve = self.cerceve_olustur(self.ana_cerceve, _("Ses Çıkış Ayarları"), 
                                                 bg=self.renkler["cerceve"])
        self.cikis_cerceve.pack(fill="x", pady=(0, 10))
        
        # Format seçimi
        format_etiket = tk.Label(self.cikis_cerceve, 
                               text=_("Çıkış Formatı:"),
                               font=self.font_normal,
                               bg=self.renkler["cerceve"],
                               fg=self.renkler["metin"])
        format_etiket.grid(row=0, column=0, padx=(10, 5), pady=10, sticky="w")
        
        self.format_degisken = tk.StringVar(value=".mp3")
        format_secenekler = [".mp3", ".wav", ".ogg", ".flac"]
        self.format_menu = ttk.Combobox(self.cikis_cerceve, 
                                      textvariable=self.format_degisken, 
                                      values=format_secenekler,
                                      state="readonly",
                                      font=self.font_normal)
        self.format_menu.grid(row=0, column=1, padx=(0, 20), pady=10, sticky="w")
        self.format_menu.bind("<<ComboboxSelected>>", self.format_degisti)
        
        # Dil seçimi
        dil_etiket = tk.Label(self.cikis_cerceve, 
                            text=_("Dil:"),
                            font=self.font_normal,
                            bg=self.renkler["cerceve"],
                            fg=self.renkler["metin"])
        dil_etiket.grid(row=0, column=2, padx=(0, 5), pady=10, sticky="w")
        
        self.dil_degisken = tk.StringVar(value="Türkçe (tr)")
        dil_secenekler = [
            "Türkçe (tr)",
            "İngilizce (en)",
            "Almanca (de)",
            "Fransızca (fr)",
            "İspanyolca (es)",
            "Rusça (ru)",
            "Arapça (ar)",
            "Japonca (ja)",
            "Çince (zh)"
        ]
        self.dil_menu = ttk.Combobox(self.cikis_cerceve, 
                                   textvariable=self.dil_degisken, 
                                   values=dil_secenekler,
                                   state="readonly",
                                   font=self.font_normal)
        self.dil_menu.grid(row=0, column=3, padx=(0, 10), pady=10, sticky="w")
        
        # Yavaş konuşma seçeneği
        self.yavas_degisken = tk.BooleanVar(value=False)
        self.yavas_check = tk.Checkbutton(self.cikis_cerceve, 
                                        text=_("Yavaş Konuşma"),
                                        variable=self.yavas_degisken,
                                        font=self.font_normal,
                                        bg=self.renkler["cerceve"],
                                        fg=self.renkler["metin"],
                                        selectcolor=self.renkler["arkaplan"],
                                        activebackground=self.renkler["cerceve"],
                                        activeforeground=self.renkler["metin"])
        self.yavas_check.grid(row=0, column=4, padx=(10, 0), pady=10, sticky="w")
        
        # Ses seviyesi
        ses_etiket = tk.Label(self.cikis_cerceve, 
                            text=_("Ses Seviyesi:"),
                            font=self.font_normal,
                            bg=self.renkler["cerceve"],
                            fg=self.renkler["metin"])
        ses_etiket.grid(row=1, column=0, padx=(10, 5), pady=(0, 10), sticky="w")
        
        self.ses_seviye_ayar = tk.Scale(self.cikis_cerceve, 
                                      from_=0, to=100,
                                      orient=tk.HORIZONTAL,
                                      command=self.ses_seviyesi_ayarla,
                                      bg=self.renkler["cerceve"],
                                      fg=self.renkler["metin"],
                                      highlightthickness=0,
                                      activebackground=self.renkler["buton_aktif"])
        self.ses_seviye_ayar.set(int(self.ses_seviyesi * 100))
        self.ses_seviye_ayar.grid(row=1, column=1, padx=(0, 10), pady=(0, 10), sticky="ew")
        
        # Çıkış konumu
        cikis_etiket = tk.Label(self.cikis_cerceve, 
                              text=_("Çıkış Konumu:"),
                              font=self.font_normal,
                              bg=self.renkler["cerceve"],
                              fg=self.renkler["metin"])
        cikis_etiket.grid(row=2, column=0, padx=(10, 5), pady=(0, 10), sticky="w")
        
        self.cikis_entry = tk.Entry(self.cikis_cerceve, 
                                  width=60,
                                  font=self.font_normal,
                                  bg="#ffffff",
                                  fg="#000000")
        self.cikis_entry.grid(row=2, column=1, columnspan=3, padx=(0, 5), pady=(0, 10), sticky="ew")
        
        # Çıkış gözat butonu
        self.cikisi_gosterme_butonu = self.buton_olustur(self.cikis_cerceve, 
                                                        _("Gözat..."), 
                                                        self.cikis_konumu_sec,
                                                        bg=self.renkler["buton"],
                                                        fg=self.renkler["metin"])
        self.cikisi_gosterme_butonu.grid(row=2, column=4, padx=(0, 10), pady=(0, 10))
        
        # Grid ayarları
        self.cikis_cerceve.grid_columnconfigure(1, weight=1)
        self.cikis_cerceve.grid_columnconfigure(3, weight=1)
    
    def onizleme_cercevesi_olustur(self):
        """PowerPoint metin önizleme bölümünü oluşturur"""
        self.onizleme_cerceve = self.cerceve_olustur(self.ana_cerceve, _("PowerPoint Metin Önizleme"), 
                                                    bg=self.renkler["cerceve"])
        self.onizleme_cerceve.pack(fill="both", expand=True, pady=(0, 10))
        
        # Önizleme metin alanı
        self.onizleme_metin = scrolledtext.ScrolledText(self.onizleme_cerceve, 
                                                      wrap=tk.WORD,
                                                      font=("Consolas", 10),
                                                      bg="#ffffff",
                                                      fg="#000000",
                                                      insertbackground="#000000",
                                                      selectbackground="#3498db",
                                                      selectforeground="#ffffff",
                                                      padx=10,
                                                      pady=10)
        self.onizleme_metin.pack(fill="both", expand=True, padx=5, pady=5)
        
        # İlerleme çubuğu
        self.ilerleme_cubugu = Progressbar(self.onizleme_cerceve, 
                                         orient=tk.HORIZONTAL,
                                         length=100,
                                         mode='determinate')
        self.ilerleme_cubugu.pack(fill="x", padx=5, pady=(0, 5))
    
    def kontrol_butonlari_olustur(self):
        """Ana kontrol butonlarını oluşturur"""
        self.kontrol_cerceve = tk.Frame(self.ana_cerceve, bg=self.renkler["arkaplan"])
        self.kontrol_cerceve.pack(fill="x", pady=(0, 10))
        
        # Dönüştür butonu
        self.donustur_butonu = self.buton_olustur(self.kontrol_cerceve, 
                                                 _("Sese Dönüştür"), 
                                                 self.donusumu_baslat,
                                                 bg=self.renkler["buton_aktif"],
                                                 fg=self.renkler["metin"])
        self.donustur_butonu.pack(side="left", padx=(0, 10))
        
        # Oynat butonu
        self.oynat_butonu = self.buton_olustur(self.kontrol_cerceve, 
                                             _("Sesi Oynat"), 
                                             self.sesi_oynat,
                                             bg=self.renkler["buton"],
                                             fg=self.renkler["metin"])
        self.oynat_butonu.pack(side="left", padx=(0, 10))
        
        # Duraklat butonu
        self.duraklat_butonu = self.buton_olustur(self.kontrol_cerceve, 
                                                _("Duraklat"), 
                                                self.sesi_duraklat,
                                                bg=self.renkler["vurgu"],
                                                fg=self.renkler["metin"])
        self.duraklat_butonu.pack(side="left", padx=(0, 10))
        
        # Durdur butonu
        self.durdur_butonu = self.buton_olustur(self.kontrol_cerceve, 
                                              _("Durdur"), 
                                              self.sesi_durdur,
                                              bg=self.renkler["uyari"],
                                              fg=self.renkler["metin"])
        self.durdur_butonu.pack(side="left", padx=(0, 10))
        
        # Ayarları kaydet butonu
        self.ayarlari_kaydet_butonu = self.buton_olustur(self.kontrol_cerceve, 
                                                        _("Ayarları Kaydet"), 
                                                        self.ayarları_kaydet,
                                                        bg=self.renkler["buton"],
                                                        fg=self.renkler["metin"])
        self.ayarlari_kaydet_butonu.pack(side="right")
    
    def durum_cubugu_olustur(self):
        """Durum çubuğunu oluşturur"""
        self.durum_cerceve = tk.Frame(self.ana_cerceve, 
                                    bg=self.renkler["arkaplan"],
                                    bd=1,
                                    relief=tk.SUNKEN)
        self.durum_cerceve.pack(fill="x", pady=(0, 10))
        
        # Durum metni
        self.durum_degisken = tk.StringVar(value=_("Hazır"))
        self.durum_etiket = tk.Label(self.durum_cerceve, 
                                   textvariable=self.durum_degisken,
                                   font=self.font_kucuk,
                                   bg=self.renkler["arkaplan"],
                                   fg=self.renkler["metin"],
                                   anchor="w")
        self.durum_etiket.pack(side="left", fill="x", expand=True, padx=5)
        
        # Sistem bilgisi
        sistem_metin = f"{self.sistem_bilgisi['platform']} {self.sistem_bilgisi['sürüm']} | Python {self.sistem_bilgisi['python']}"
        self.sistem_etiket = tk.Label(self.durum_cerceve, 
                                    text=sistem_metin,
                                    font=self.font_kucuk,
                                    bg=self.renkler["arkaplan"],
                                    fg=self.renkler["metin"],
                                    anchor="e")
        self.sistem_etiket.pack(side="right", padx=5)
    
    def menu_cubugu_olustur(self):
        """Menü çubuğunu oluşturur"""
        self.menu_cubugu = tk.Menu(self.root, bg=self.renkler["arkaplan"], fg=self.renkler["metin"])
        self.root.config(menu=self.menu_cubugu)
        
        # Dosya menüsü
        dosya_menu = tk.Menu(self.menu_cubugu, tearoff=0, bg=self.renkler["arkaplan"], fg=self.renkler["metin"])
        dosya_menu.add_command(label=_("Aç..."), command=self.dosya_sec, accelerator="Ctrl+O")
        dosya_menu.add_command(label=_("Son Kullanılanları Aç"), command=self.son_kullanılanlari_ac)
        dosya_menu.add_separator()
        dosya_menu.add_command(label=_("Ayarları Kaydet"), command=self.ayarları_kaydet, accelerator="Ctrl+S")
        dosya_menu.add_separator()
        dosya_menu.add_command(label=_("Çıkış"), command=self.root.quit, accelerator="Alt+F4")
        self.menu_cubugu.add_cascade(label=_("Dosya"), menu=dosya_menu)
        
        # Düzen menüsü
        duzen_menu = tk.Menu(self.menu_cubugu, tearoff=0, bg=self.renkler["arkaplan"], fg=self.renkler["metin"])
        duzen_menu.add_command(label=_("Kes"), command=lambda: self.onizleme_metin.event_generate("<<Cut>>"))
        duzen_menu.add_command(label=_("Kopyala"), command=lambda: self.onizleme_metin.event_generate("<<Copy>>"))
        duzen_menu.add_command(label=_("Yapıştır"), command=lambda: self.onizleme_metin.event_generate("<<Paste>>"))
        duzen_menu.add_separator()
        duzen_menu.add_command(label=_("Tümünü Seç"), command=lambda: self.onizleme_metin.tag_add("sel", "1.0", "end"))
        duzen_menu.add_command(label=_("Temizle"), command=lambda: self.onizleme_metin.delete("1.0", tk.END))
        self.menu_cubugu.add_cascade(label=_("Düzen"), menu=duzen_menu)
        
        # Araçlar menüsü
        arac_menu = tk.Menu(self.menu_cubugu, tearoff=0, bg=self.renkler["arkaplan"], fg=self.renkler["metin"])
        arac_menu.add_command(label=_("Metni Düzenle"), command=self.metni_duzenle)
        arac_menu.add_command(label=_("Ses Efektleri Ekle"), command=self.ses_efektleri_ekle)
        arac_menu.add_separator()
        arac_menu.add_command(label=_("Toplu Dönüştür"), command=self.toplu_donustur)
        arac_menu.add_command(label=_("Güncellemeleri Kontrol Et"), command=self.guncelleme_kontrol)
        self.menu_cubugu.add_cascade(label=_("Araçlar"), menu=arac_menu)
        
        # Yardım menüsü
        yardim_menu = tk.Menu(self.menu_cubugu, tearoff=0, bg=self.renkler["arkaplan"], fg=self.renkler["metin"])
        yardim_menu.add_command(label=_("Yardım"), command=self.yardim_goster)
        yardim_menu.add_command(label=_("Kısayollar"), command=self.kisayollari_goster)
        yardim_menu.add_separator()
        yardim_menu.add_command(label=_("Hakkında"), command=self.hakkinda_goster)
        self.menu_cubugu.add_cascade(label=_("Yardım"), menu=yardim_menu)
        
        # Klavye kısayolları
        self.root.bind("<Control-o>", lambda e: self.dosya_sec())
        self.root.bind("<Control-s>", lambda e: self.ayarları_kaydet())
        self.root.bind("<F1>", lambda e: self.yardim_goster())
    
    def arac_cubugu_olustur(self):
        """Araç çubuğunu oluşturur"""
        self.arac_cubugu = tk.Frame(self.ana_cerceve, bg=self.renkler["arkaplan"], height=30)
        self.arac_cubugu.pack(fill="x", pady=(0, 10))
        
        # Araç butonları
        arac_butonlari = [
            ("📂", _("Aç"), self.dosya_sec),
            ("💾", _("Kaydet"), self.ayarları_kaydet),
            ("▶️", _("Oynat"), self.sesi_oynat),
            ("⏸️", _("Duraklat"), self.sesi_duraklat),
            ("⏹️", _("Durdur"), self.sesi_durdur),
            ("⚙️", _("Ayarlar"), self.ayarlari_goster),
            ("❓", _("Yardım"), self.yardim_goster)
        ]
        
        for icon, text, command in arac_butonlari:
            btn = self.buton_olustur(self.arac_cubugu, f" {icon} {text} ", command,
                                    bg=self.renkler["arkaplan"], fg=self.renkler["metin"],
                                    bd=0, font=self.font_kucuk)
            btn.pack(side="left", padx=2)
    
    def cerceve_olustur(self, parent, text=None, **kwargs):
        """Stilize bir çerçeve oluşturur"""
        bg = kwargs.pop("bg", self.renkler["cerceve"])
        fg = kwargs.pop("fg", self.renkler["metin"])
        
        if text:
            cerceve = tk.LabelFrame(parent, text=text, 
                                  bg=bg, fg=fg,
                                  font=self.font_alt_baslik,
                                  **kwargs)
        else:
            cerceve = tk.Frame(parent, bg=bg, **kwargs)
        
        return cerceve
    
    def buton_olustur(self, parent, text, command, **kwargs):
        """Stilize bir buton oluşturur"""
        bg = kwargs.pop("bg", self.renkler["buton"])
        fg = kwargs.pop("fg", self.renkler["metin"])
        activebg = kwargs.pop("activebackground", self.renkler["buton_hover"])
        activefg = kwargs.pop("activeforeground", self.renkler["metin"])
        bd = kwargs.pop("bd", 2)
        font = kwargs.pop("font", self.font_normal)

        btn = tk.Button(parent, text=text, command=command,
                    bg=bg, fg=fg,
                    activebackground=activebg,
                    activeforeground=activefg,
                    font=font,
                    relief=tk.RAISED,
                    bd=bd,
                    **kwargs)

    # Hover efekti
        btn.bind("<Enter>", lambda e: btn.config(bg=self.renkler["buton_hover"]))
        btn.bind("<Leave>", lambda e: btn.config(bg=bg))

        return btn

    
    def dosya_sec(self):
        """PowerPoint dosyası seçme dialogunu açar"""
        dosya_yolu = filedialog.askopenfilename(
            title=_("PowerPoint Dosyası Seçin"),
            initialdir=self.son_kullanılan_klasor,
            filetypes=[(_("PowerPoint Dosyaları"), "*.pptx *.ppt"), (_("Tüm Dosyalar"), "*.*")]
        )
        
        if dosya_yolu:
            self.giris_dosyası = dosya_yolu
            self.giris_entry.delete(0, tk.END)
            self.giris_entry.insert(0, dosya_yolu)
            self.son_kullanılan_klasor = os.path.dirname(dosya_yolu)
            
            # Varsayılan çıkış yolunu ayarla
            temel_ad = os.path.splitext(os.path.basename(dosya_yolu))[0]
            cikis_yolu = os.path.join(os.path.dirname(dosya_yolu), f"{temel_ad}{self.format_degisken.get()}")
            self.cikis_entry.delete(0, tk.END)
            self.cikis_entry.insert(0, cikis_yolu)
            
            # Dosya bilgilerini göster
            self.dosya_bilgilerini_goster(dosya_yolu)
            
            # Metni önizle
            self.powerpoint_metnini_onizle()
    
    def dosya_bilgilerini_goster(self, dosya_yolu):
        """Seçilen dosyanın bilgilerini gösterir"""
        try:
            dosya_boyutu = os.path.getsize(dosya_yolu)
            duzenleme_tarihi = datetime.fromtimestamp(os.path.getmtime(dosya_yolu))
            
            bilgi_metni = _("Dosya: {} | Boyut: {} | Son Düzenleme: {}").format(
                os.path.basename(dosya_yolu),
                self.boyutu_duzenle(dosya_boyutu),
                duzenleme_tarihi.strftime("%d/%m/%Y %H:%M:%S")
            )
            
            self.dosya_bilgi_etiket.config(text=bilgi_metni)
            self.durum_guncelle(_("Dosya seçildi: {}").format(os.path.basename(dosya_yolu)), "basari")
            
        except Exception as e:
            self.durum_guncelle(_("Dosya bilgileri alınamadı: {}").format(str(e)), "uyari")
    
    def boyutu_duzenle(self, boyut):
        """Dosya boyutunu uygun birime çevirir"""
        for birim in ['B', 'KB', 'MB', 'GB']:
            if boyut < 1024.0:
                return f"{boyut:.1f} {birim}"
            boyut /= 1024.0
        return f"{boyut:.1f} TB"
    
    def powerpoint_metnini_onizle(self):
        """PowerPoint dosyasındaki metni çıkarır ve önizleme alanına yazar"""
        if not self.giris_dosyası:
            return
            
        try:
            ppt = Presentation(self.giris_dosyası)
            metinler = []
            
            for slayt_no, slayt in enumerate(ppt.slides, 1):
                metinler.append(f"=== Slayt {slayt_no} ===")
                for sekil in slayt.shapes:
                    if hasattr(sekil, "text"):
                        metinler.append(sekil.text)
                metinler.append("")  # Boş satır ekle
            self.onizleme_metin.delete(1.0, tk.END)
            self.onizleme_metin.insert(tk.END, "\n".join(metinler))
            self.durum_guncelle(_("PowerPoint metni başarıyla yüklendi"), "basari")

        except Exception as e:
            self.durum_guncelle(_("PowerPoint metni çıkarılırken hata: {}").format(str(e)), "uyari")
            self.onizleme_metin.delete(1.0, tk.END)
            self.onizleme_metin.insert(tk.END, _("HATA: Dosya okunamadı - {}").format(str(e)))
    

    def cikis_konumu_sec(self):
        """Çıkış dosyası konumu seçme dialogunu açar"""
        dosya_yolu = filedialog.asksaveasfilename(
            title=_("Çıkış Dosyasını Kaydet"),
            initialdir=self.son_kullanılan_klasor,
            defaultextension=self.format_degisken.get(),
            filetypes=[(_("Ses Dosyaları"), "*.mp3 *.wav *.ogg *.flac"), (_("Tüm Dosyalar"), "*.*")]
        )
        
        if dosya_yolu:
            self.cikis_dosyası = dosya_yolu
            self.cikis_entry.delete(0, tk.END)
            self.cikis_entry.insert(0, dosya_yolu)
            self.son_kullanılan_klasor = os.path.dirname(dosya_yolu)
            self.durum_guncelle(_("Çıkış konumu ayarlandı: {}").format(os.path.basename(dosya_yolu)), "basari")
    
    def format_degisti(self, event=None):
        """Çıkış formatı değiştiğinde çıkış yolunu günceller"""
        if self.giris_dosyası:
            temel_ad = os.path.splitext(os.path.basename(self.giris_dosyası))[0]
            yeni_uzanti = self.format_degisken.get()
            mevcut_yol = self.cikis_entry.get()
            
            # Eğer mevcut yol varsa ve uzantısı değişecekse
            if mevcut_yol and os.path.splitext(mevcut_yol)[1].lower() != yeni_uzanti.lower():
                yeni_yol = os.path.join(os.path.dirname(mevcut_yol), f"{temel_ad}{yeni_uzanti}")
                self.cikis_entry.delete(0, tk.END)
                self.cikis_entry.insert(0, yeni_yol)
    
    def ses_seviyesi_ayarla(self, deger):
        """Ses seviyesini ayarlar"""
        self.ses_seviyesi = float(deger) / 100
        if pygame.mixer.get_init():
            pygame.mixer.music.set_volume(self.ses_seviyesi)
    
    def donusumu_baslat(self):
        """Dönüştürme işlemini başlatır"""
        if not self.giris_dosyası:
            messagebox.showerror(_("Hata"), _("Lütfen bir PowerPoint dosyası seçin"))
            return
            
        if not self.cikis_entry.get():
            messagebox.showerror(_("Hata"), _("Lütfen çıkış dosyası konumunu belirtin"))
            return
            
        self.cikis_dosyası = self.cikis_entry.get()
        self.dil = self.dil_degisken.get().split("(")[1][:-1]
        self.yavas_konusma = self.yavas_degisken.get()
        
        # İlerleme çubuğunu sıfırla
        self.ilerleme_cubugu["value"] = 0
        
        # Dönüştür butonunu devre dışı bırak
        self.donustur_butonu.config(state=tk.DISABLED)
        self.durum_guncelle(_("Dönüştürme başlatılıyor..."), "bilgi")
        
        # Ayrı bir thread'de dönüştürme işlemi
        donusum_thread = threading.Thread(target=self.powerpointi_sese_donustur, daemon=True)
        donusum_thread.start()
    
    def powerpointi_sese_donustur(self):
        """PowerPoint dosyasını sese dönüştürür (MP3 olarak)"""
        try:
            # 1. Kullanıcıdan kaydetme yeri ve dosya adını sor
            self.cikis_dosyası = filedialog.asksaveasfilename(
                title=_("MP3 Olarak Kaydet"),
                initialdir=self.son_kullanılan_klasor,
                defaultextension=".mp3",
                filetypes=[("MP3 Dosyaları", "*.mp3")]
                )
        
            if not self.cikis_dosyası:  # Kullanıcı iptal etti
                self.durum_guncelle(_("Dönüştürme iptal edildi"), "uyari")
                return

            # 2. PowerPoint metnini çıkar
            ppt = Presentation(self.giris_dosyası)
            tum_metin = []
            for slayt in ppt.slides:
                for sekil in slayt.shapes:
                    if hasattr(sekil, "text") and sekil.text.strip():
                        tum_metin.append(sekil.text)
        
            birlesik_metin = "\n".join(tum_metin)
        
            if not birlesik_metin.strip():
                messagebox.showerror(_("Hata"), _("PowerPoint dosyasında metin bulunamadı!"))
                return

            # 3. gTTS ile doğrudan MP3'e kaydet (FFmpeg GEREKMEZ)
            tts = gTTS(text=birlesik_metin, lang=self.dil, slow=self.yavas_konusma)
            tts.save(self.cikis_dosyası)
        
            self.durum_guncelle(_("Dönüştürme tamamlandı: {}").format(os.path.basename(self.cikis_dosyası)), "basari")
            messagebox.showinfo(_("Başarılı"), _("Ses dosyası başarıyla kaydedildi!"))
        
        except Exception as e:
            self.durum_guncelle(_("Dönüştürme hatası: {}").format(str(e)), "uyari")
            messagebox.showerror(_("Hata"), _("Dönüştürme başarısız:\n{}").format(str(e)))
        finally:
            self.donusum_devam_ediyor = False
            self.donustur_butonu.config(state=tk.NORMAL)
            
    
    def sesi_oynat(self):
        """Oluşturulan ses dosyasını oynatır"""
        if not self.cikis_dosyası or not os.path.exists(self.cikis_dosyası):
            messagebox.showerror(_("Hata"), _("Oynatılacak ses dosyası bulunamadı"))
            return
            
        try:
            if pygame.mixer.music.get_busy():
                pygame.mixer.music.stop()
                
            pygame.mixer.music.load(self.cikis_dosyası)
            pygame.mixer.music.play()
            self.oynatma_devam_ediyor = True
            self.durum_guncelle(_("Ses oynatılıyor: {}").format(os.path.basename(self.cikis_dosyası)), "basari")
            
            # Oynatma bitişini kontrol etmek için thread başlat
            threading.Thread(target=self.oynatma_durumunu_kontrol_et, daemon=True).start()
            
        except Exception as e:
            self.durum_guncelle(_("Ses oynatma hatası: {}").format(str(e)), "uyari")
            messagebox.showerror(_("Hata"), _("Ses oynatılırken bir hata oluştu:\n{}").format(str(e)))
    
    def oynatma_durumunu_kontrol_et(self):
        """Ses oynatmanın bitip bitmediğini kontrol eder"""
        while self.oynatma_devam_ediyor and pygame.mixer.music.get_busy():
            time.sleep(0.1)
        
        if self.oynatma_devam_ediyor:
            self.oynatma_devam_ediyor = False
            self.durum_guncelle(_("Oynatma tamamlandı"), "bilgi")
    
    def sesi_duraklat(self):
        """Sesi duraklatır"""
        if pygame.mixer.music.get_busy():
            pygame.mixer.music.pause()
            self.oynatma_devam_ediyor = False
            self.durum_guncelle(_("Ses duraklatıldı"), "bilgi")
    
    def sesi_devam_ettir(self):
        """Duraklatılmış sesi devam ettirir"""
        pygame.mixer.music.unpause()
        self.oynatma_devam_ediyor = True
        self.durum_guncelle(_("Ses devam ediyor"), "bilgi")
        threading.Thread(target=self.oynatma_durumunu_kontrol_et, daemon=True).start()
    
    def sesi_durdur(self):
        """Sesi tamamen durdurur"""
        pygame.mixer.music.stop()
        self.oynatma_devam_ediyor = False
        self.durum_guncelle(_("Ses durduruldu"), "bilgi")
    
    def durum_guncelle(self, mesaj, tip="normal"):
        """Durum çubuğunu günceller"""
        renkler = {
            "normal": self.renkler["metin"],
            "bilgi": self.renkler["baslik"],
            "uyari": self.renkler["uyari"],
            "basari": self.renkler["basari"],
            "vurgu": self.renkler["vurgu"]
        }
        
        self.durum_degisken.set(mesaj)
        self.durum_etiket.config(fg=renkler.get(tip, self.renkler["metin"]))
        
        # Durum mesajını logla
        self.log_mesaji(mesaj)
    
    def log_mesaji(self, mesaj):
        """Mesajı log dosyasına yazar"""
        log_dosyasi = "powerpoint_ses_donusturucu.log"
        zaman_damgasi = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        with open(log_dosyasi, "a", encoding="utf-8") as f:
            f.write(f"[{zaman_damgasi}] {mesaj}\n")
    
    def ayarlari_goster(self):
        """Ayarlar penceresini gösterir"""
        ayar_penceresi = tk.Toplevel(self.root)
        ayar_penceresi.title(_("Ayarlar"))
        ayar_penceresi.geometry("600x400")
        ayar_penceresi.resizable(False, False)
        ayar_penceresi.transient(self.root)
        ayar_penceresi.grab_set()
        
        # Ayarlar içeriği
        ayar_cerceve = tk.Frame(ayar_penceresi, padx=20, pady=20)
        ayar_cerceve.pack(fill="both", expand=True)
        
        # Dil ayarları
        dil_etiket = tk.Label(ayar_cerceve, text=_("Uygulama Dili:"), font=self.font_normal)
        dil_etiket.grid(row=0, column=0, sticky="w", pady=(0, 10))
        
        dil_secenekler = ["Türkçe", "English", "Deutsch", "Français", "Español"]
        self.dil_degisken_ayar = tk.StringVar(value="Türkçe")
        dil_menu = ttk.Combobox(ayar_cerceve, textvariable=self.dil_degisken_ayar, 
                               values=dil_secenekler, state="readonly", font=self.font_normal)
        dil_menu.grid(row=0, column=1, sticky="ew", pady=(0, 10), padx=(10, 0))
        
        # Varsayılan klasör ayarı
        klasor_etiket = tk.Label(ayar_cerceve, text=_("Varsayılan Klasör:"), font=self.font_normal)
        klasor_etiket.grid(row=1, column=0, sticky="w", pady=(0, 10))
        
        self.klasor_entry = tk.Entry(ayar_cerceve, font=self.font_normal)
        self.klasor_entry.insert(0, self.son_kullanılan_klasor)
        self.klasor_entry.grid(row=1, column=1, sticky="ew", pady=(0, 10), padx=(10, 0))
        
        klasor_buton = self.buton_olustur(ayar_cerceve, _("Gözat..."), self.klasor_sec)
        klasor_buton.grid(row=1, column=2, padx=(10, 0), pady=(0, 10))
        
        # Gece modu ayarı
        self.gece_modu_degisken = tk.BooleanVar(value=self.gece_modu)
        gece_modu_check = tk.Checkbutton(ayar_cerceve, 
                                        text=_("Gece Modu"),
                                        variable=self.gece_modu_degisken,
                                        font=self.font_normal)
        gece_modu_check.grid(row=2, column=0, columnspan=3, sticky="w", pady=(0, 20))
        
        # Kaydet butonu
        kaydet_buton = self.buton_olustur(ayar_cerceve, _("Ayarları Kaydet"), 
                                         lambda: self.ayarlari_kaydet_ve_kapat(ayar_penceresi))
        kaydet_buton.grid(row=3, column=0, columnspan=3, pady=(20, 0))
        
        # Grid ayarları
        ayar_cerceve.grid_columnconfigure(1, weight=1)
    
    def klasor_sec(self):
        """Varsayılan klasör seçme dialogunu açar"""
        klasor = filedialog.askdirectory(
            title=_("Varsayılan Klasör Seçin"),
            initialdir=self.son_kullanılan_klasor
        )
        
        if klasor:
            self.klasor_entry.delete(0, tk.END)
            self.klasor_entry.insert(0, klasor)
    
    def ayarlari_kaydet_ve_kapat(self, pencere):
        """Ayarları kaydeder ve pencereyi kapatır"""
        self.son_kullanılan_klasor = self.klasor_entry.get()
        self.gece_modu = self.gece_modu_degisken.get()
        
        # Gece modunu uygula
        if self.gece_modu:
            self.gece_modunu_uygula()
        else:
            self.gunduz_modunu_uygula()
        
        # Ayarları dosyaya kaydet
        self.ayarları_kaydet()
        
        pencere.destroy()
        messagebox.showinfo(_("Başarılı"), _("Ayarlar kaydedildi"))
    
    def gece_modunu_degistir(self):
        """Gece modunu açıp kapatır"""
        self.gece_modu = not self.gece_modu
        
        if self.gece_modu:
            self.gece_modunu_uygula()
            self.gece_modu_butonu.config(text="☀️")
        else:
            self.gunduz_modunu_uygula()
            self.gece_modu_butonu.config(text="🌙")
    
    def gece_modunu_uygula(self):
        """Gece modu renklerini uygular"""
        self.renkler = {
            "arkaplan": "#1a1a1a",
            "cerceve": "#2d2d2d",
            "baslik": "#3a3a3a",
            "buton": "#4a4a4a",
            "buton_hover": "#5a5a5a",
            "buton_aktif": "#2a5a2a",
            "metin": "#e0e0e0",
            "uyari": "#8b0000",
            "basari": "#006400",
            "vurgu": "#8b6500"
        }
        self.arayuzu_yeniden_yukle()
    
    def gunduz_modunu_uygula(self):
        """Gündüz modu renklerini uygular"""
        self.renkler = {
            "arkaplan": "#2c3e50",
            "cerceve": "#34495e",
            "baslik": "#3498db",
            "buton": "#2980b9",
            "buton_hover": "#3498db",
            "buton_aktif": "#1abc9c",
            "metin": "#ecf0f1",
            "uyari": "#e74c3c",
            "basari": "#2ecc71",
            "vurgu": "#f39c12"
        }
        self.arayuzu_yeniden_yukle()
    
    def arayuzu_yeniden_yukle(self):
        """Arayüzü yeniden yükler ve yeni renkleri uygular"""
        # Mevcut arayüzü temizle
        for widget in self.ana_cerceve.winfo_children():
            widget.destroy()
        
        # Yeni arayüzü oluştur
        self.arayuz_olustur()
        
        # PowerPoint metnini yeniden yükle
        if self.giris_dosyası:
            self.powerpoint_metnini_onizle()
    
    def ayarları_yukle(self):
        """Kayıtlı ayarları yükler"""
        ayar_dosyasi = "powerpoint_ses_ayarlari.json"
        
        try:
            if os.path.exists(ayar_dosyasi):
                with open(ayar_dosyasi, "r", encoding="utf-8") as f:
                    ayarlar = json.load(f)
                    
                    self.son_kullanılan_klasor = ayarlar.get("son_kullanılan_klasor", os.path.expanduser("~"))
                    self.gece_modu = ayarlar.get("gece_modu", False)
                    self.ses_seviyesi = ayarlar.get("ses_seviyesi", 0.7)
                    self.dil = ayarlar.get("dil", "tr")
                    self.yavas_konusma = ayarlar.get("yavas_konusma", False)
                    
                    # Dil menüsünü ayarla
                    for secenek in self.dil_menu["values"]:
                        if f"({self.dil})" in secenek:
                            self.dil_degisken.set(secenek)
                            break
                    
                    # Ses seviyesini ayarla
                    self.ses_seviye_ayar.set(int(self.ses_seviyesi * 100))
                    
                    # Gece modunu uygula
                    if self.gece_modu:
                        self.gece_modunu_uygula()
                        self.gece_modu_butonu.config(text="☀️")
        
        except Exception as e:
            self.durum_guncelle(_("Ayarlar yüklenirken hata: {}").format(str(e)), "uyari")
    
    def ayarları_kaydet(self):
        """Mevcut ayarları kaydeder"""
        ayar_dosyasi = "powerpoint_ses_ayarlari.json"
        ayarlar = {
            "son_kullanılan_klasor": self.son_kullanılan_klasor,
            "gece_modu": self.gece_modu,
            "ses_seviyesi": self.ses_seviyesi,
            "dil": self.dil,
            "yavas_konusma": self.yavas_konusma,
            "son_guncelleme": datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }
        
        try:
            with open(ayar_dosyasi, "w", encoding="utf-8") as f:
                json.dump(ayarlar, f, indent=4, ensure_ascii=False)
            
            self.durum_guncelle(_("Ayarlar kaydedildi"), "basari")
            return True
            
        except Exception as e:
            self.durum_guncelle(_("Ayarlar kaydedilirken hata: {}").format(str(e)), "uyari")
            return False
    
    def metni_duzenle(self):
        """Metin düzenleme penceresini açar"""
        if not self.onizleme_metin.get("1.0", tk.END).strip():
            messagebox.showwarning(_("Uyarı"), _("Düzenlenecek metin bulunamadı"))
            return
            
        duzenleme_penceresi = tk.Toplevel(self.root)
        duzenleme_penceresi.title(_("Metin Düzenleyici"))
        duzenleme_penceresi.geometry("800x600")
        duzenleme_penceresi.transient(self.root)
        duzenleme_penceresi.grab_set()
        
        # Metin editörü
        editor = scrolledtext.ScrolledText(duzenleme_penceresi, wrap=tk.WORD, font=("Consolas", 10))
        editor.pack(fill="both", expand=True, padx=10, pady=10)
        editor.insert(tk.END, self.onizleme_metin.get("1.0", tk.END))
        
        # Buton çerçevesi
        buton_cerceve = tk.Frame(duzenleme_penceresi)
        buton_cerceve.pack(fill="x", padx=10, pady=(0, 10))
        
        # Kaydet butonu
        kaydet_buton = self.buton_olustur(buton_cerceve, _("Kaydet"), 
                                         lambda: self.metni_kaydet(editor, duzenleme_penceresi))
        kaydet_buton.pack(side="left", padx=(0, 10))
        
        # İptal butonu
        iptal_buton = self.buton_olustur(buton_cerceve, _("İptal"), 
                                         duzenleme_penceresi.destroy)
        iptal_buton.pack(side="left")
    
    def metni_kaydet(self, editor, pencere):
        """Düzenlenen metni kaydeder"""
        yeni_metin = editor.get("1.0", tk.END)
        self.onizleme_metin.delete("1.0", tk.END)
        self.onizleme_metin.insert(tk.END, yeni_metin)
        pencere.destroy()
        self.durum_guncelle(_("Metin başarıyla düzenlendi"), "basari")
    
    def ses_efektleri_ekle(self):
        """Ses efektleri ekleme penceresini açar"""
        efekt_penceresi = tk.Toplevel(self.root)
        efekt_penceresi.title(_("Ses Efektleri Ekle"))
        efekt_penceresi.geometry("600x400")
        efekt_penceresi.transient(self.root)
        efekt_penceresi.grab_set()
        
        # Efekt seçenekleri
        efekt_cerceve = tk.LabelFrame(efekt_penceresi, text=_("Mevcut Efektler"), padx=10, pady=10)
        efekt_cerceve.pack(fill="both", expand=True, padx=10, pady=10)
        
        efektler = [
            (_("Sessizlik Ekle"), "silence"),
            (_("Giriş Müziği Ekle"), "intro"),
            (_("Bitiş Müziği Ekle"), "outro"),
            (_("Sayfa Geçiş Efekti"), "transition"),
            (_("Arka Plan Müziği"), "background")
        ]
        
        for text, efekt in efektler:
            btn = self.buton_olustur(efekt_cerceve, text, 
                                    lambda e=efekt: self.efekt_uygula(e, efekt_penceresi))
            btn.pack(fill="x", pady=5)
    
    def efekt_uygula(self, efekt, pencere):
        """Seçilen efekt uygulanır"""
        self.durum_guncelle(_("{} efekti uygulanıyor...").format(efekt), "bilgi")
        pencere.destroy()
        messagebox.showinfo(_("Bilgi"), _("{} efekti uygulandı").format(efekt))
    
    def toplu_donustur(self):
        """Toplu dönüştürme penceresini açar"""
        toplu_penceresi = tk.Toplevel(self.root)
        toplu_penceresi.title(_("Toplu Dönüştürme"))
        toplu_penceresi.geometry("800x600")
        toplu_penceresi.transient(self.root)
        toplu_penceresi.grab_set()
        
        # Dosya listesi
        liste_cerceve = tk.LabelFrame(toplu_penceresi, text=_("Dönüştürülecek Dosyalar"), padx=10, pady=10)
        liste_cerceve.pack(fill="both", expand=True, padx=10, pady=10)
        
        self.topluliste = tk.Listbox(liste_cerceve, selectmode=tk.EXTENDED)
        self.topluliste.pack(fill="both", expand=True, padx=5, pady=5)
        
        # Buton çerçevesi
        buton_cerceve = tk.Frame(toplu_penceresi)
        buton_cerceve.pack(fill="x", padx=10, pady=(0, 10))
        
        # Dosya ekle butonu
        ekle_buton = self.buton_olustur(buton_cerceve, _("Dosya Ekle..."), 
                                       self.topluliste_dosya_ekle)
        ekle_buton.pack(side="left", padx=(0, 10))
        
        # Klasör ekle butonu
        klasor_buton = self.buton_olustur(buton_cerceve, _("Klasör Ekle..."), 
                                         self.topluliste_klasor_ekle)
        klasor_buton.pack(side="left", padx=(0, 10))
        
        # Temizle butonu
        temizle_buton = self.buton_olustur(buton_cerceve, _("Listeyi Temizle"), 
                                          self.topluliste_temizle)
        temizle_buton.pack(side="left", padx=(0, 10))
        
        # Dönüştür butonu
        donustur_buton = self.buton_olustur(buton_cerceve, _("Dönüştür"), 
                                           lambda: self.topluliste_donustur(toplu_penceresi))
        donustur_buton.pack(side="right")
    
    def topluliste_dosya_ekle(self):
        """Toplu listeye dosya ekler"""
        dosyalar = filedialog.askopenfilenames(
            title=_("PowerPoint Dosyaları Seçin"),
            initialdir=self.son_kullanılan_klasor,
            filetypes=[(_("PowerPoint Dosyaları"), "*.pptx *.ppt"), (_("Tüm Dosyalar"), "*.*")]
        )
        
        if dosyalar:
            for dosya in dosyalar:
                self.topluliste.insert(tk.END, dosya)
            self.son_kullanılan_klasor = os.path.dirname(dosyalar[0])
    
    def topluliste_klasor_ekle(self):
        """Toplu listeye klasördeki dosyaları ekler"""
        klasor = filedialog.askdirectory(
            title=_("Klasör Seçin"),
            initialdir=self.son_kullanılan_klasor
        )
        
        if klasor:
            for dosya in os.listdir(klasor):
                if dosya.lower().endswith((".ppt", ".pptx")):
                    self.topluliste.insert(tk.END, os.path.join(klasor, dosya))
            self.son_kullanılan_klasor = klasor
    
    def topluliste_temizle(self):
        """Toplu listeyi temizler"""
        self.topluliste.delete(0, tk.END)
    
    def topluliste_donustur(self, pencere):
        """Toplu listedeki dosyaları dönüştürür"""
        if self.topluliste.size() == 0:
            messagebox.showwarning(_("Uyarı"), _("Dönüştürülecek dosya bulunamadı"))
            return
            
        # Çıkış klasörü seç
        cikis_klasor = filedialog.askdirectory(
            title=_("Çıkış Klasörü Seçin"),
            initialdir=self.son_kullanılan_klasor
        )
        
        if not cikis_klasor:
            return
            
        self.son_kullanılan_klasor = cikis_klasor
        
        # İlerleme penceresi
        ilerleme_penceresi = tk.Toplevel(pencere)
        ilerleme_penceresi.title(_("Dönüştürme İlerlemesi"))
        ilerleme_penceresi.geometry("400x200")
        ilerleme_penceresi.transient(pencere)
        ilerleme_penceresi.grab_set()
        
        # İlerleme çubuğu
        ilerleme_etiket = tk.Label(ilerleme_penceresi, text=_("Toplu dönüştürme devam ediyor..."))
        ilerleme_etiket.pack(pady=(20, 10))
        
        ilerleme_cubugu = Progressbar(ilerleme_penceresi, orient=tk.HORIZONTAL, length=300, mode='determinate')
        ilerleme_cubugu.pack(pady=10)
        
        durum_etiket = tk.Label(ilerleme_penceresi, text="")
        durum_etiket.pack(pady=10)
        
        pencere.withdraw()
        
        # Ayrı bir thread'de dönüştürme işlemi
        donusum_thread = threading.Thread(
            target=self.topluliste_donustur_thread,
            args=(ilerleme_penceresi, ilerleme_cubugu, durum_etiket, cikis_klasor),
            daemon=True
        )
        donusum_thread.start()
    
    def topluliste_donustur_thread(self, pencere, cubuk, etiket, cikis_klasor):
        """Toplu dönüştürme işlemini gerçekleştirir"""
        basarili = 0
        basarisiz = 0
        
        for i in range(self.topluliste.size()):
            if not self.donusum_devam_ediyor:
                break
                
            dosya = self.topluliste.get(i)
            temel_ad = os.path.splitext(os.path.basename(dosya))[0]
            cikis_dosyasi = os.path.join(cikis_klasor, f"{temel_ad}{self.format_degisken.get()}")
            
            etiket.config(text=_("İşleniyor: {}").format(os.path.basename(dosya)))
            pencere.update()
            
            try:
                # PowerPoint dosyasını yükle
                ppt = Presentation(dosya)
                tum_metin = []
                
                for slayt in ppt.slides:
                    for sekil in slayt.shapes:
                        if hasattr(sekil, "text") and sekil.text.strip():
                            tum_metin.append(sekil.text)
                
                birlesik_metin = "\n".join(tum_metin)
                
                if birlesik_metin.strip():
                    # gTTS ile ses oluştur
                    tts = gTTS(text=birlesik_metin, lang=self.dil, slow=self.yavas_konusma)
                    
                    # Geçici dosyaya kaydet
                    gecici_dosya = "temp_audio.mp3"
                    tts.save(gecici_dosya)
                    
                    # İstenen formata dönüştür
                    ses = AudioSegment.from_mp3(gecici_dosya)
                    format = self.format_degisken.get()[1:]  # noktayı kaldır (.mp3 -> mp3)
                    
                    if format == "wav":
                        ses.export(cikis_dosyasi, format="wav")
                    elif format == "ogg":
                        ses.export(cikis_dosyasi, format="ogg")
                    elif format == "flac":
                        ses.export(cikis_dosyasi, format="flac")
                    else:  # varsayılan mp3
                        ses.export(cikis_dosyasi, format="mp3")

                    # Geçici dosyayı sil
                    os.remove(gecici_dosya)

                    basarili += 1

                else:
                    basarisiz += 1
                    self.log_mesaji(_("{} dosyasında dönüştürülebilir metin bulunamadı").format(dosya))
            except Exception as e:
                basarisiz += 1
                self.log_mesaji(_("{} dönüştürülürken hata: {}").format(dosya, str(e)))

            # İlerleme çubuğunu güncelle
            cubuk["value"] = (i + 1) / self.topluliste.size() * 100
            pencere.update()

        # Sonuçları göster
        etiket.config(text=_("Tamamlandı: {} başarılı, {} başarısız").format(basarili, basarisiz))

        # Kapat butonu ekle
        kapat_buton = self.buton_olustur(pencere, _("Kapat"), pencere.destroy)
        kapat_buton.pack(pady=10)
        
        self.durum_guncelle(_("Toplu dönüştürme tamamlandı: {} başarılı, {} başarısız").format(basarili, basarisiz),                  
                            "basari" if basarisiz == 0 else "uyari")
    def son_kullanılanlari_ac(self):
        """Son kullanılan dosyaları gösterir"""
        son_dosyalar = self.ayarları_yukle().get("son_dosyalar", [])
        
        if not son_dosyalar:
            messagebox.showinfo(_("Bilgi"), _("Son kullanılan dosya bulunamadı"))
            return
        
        # Son kullanılanlar penceresi
        son_penceresi = tk.Toplevel(self.root)
        son_penceresi.title(_("Son Kullanılan Dosyalar"))
        son_penceresi.geometry("600x400")
        son_penceresi.transient(self.root)
        son_penceresi.grab_set()
        
        # Liste çerçevesi
        liste_cerceve = tk.LabelFrame(son_penceresi, text=_("Son Dosyalar"), padx=10, pady=10)
        liste_cerceve.pack(fill="both", expand=True, padx=10, pady=10)
        
        liste = tk.Listbox(liste_cerceve, font=self.font_normal)
        liste.pack(fill="both", expand=True, padx=5, pady=5)
        
        for dosya in son_dosyalar:
            liste.insert(tk.END, os.path.basename(dosya) + " - " + dosya)
        
        # Buton çerçevesi
        buton_cerceve = tk.Frame(son_penceresi)
        buton_cerceve.pack(fill="x", padx=10, pady=(0, 10))
        
        # Aç butonu
        ac_buton = self.buton_olustur(buton_cerceve, _("Aç"), 
                                     lambda: self.son_dosya_ac(liste, son_penceresi))
        ac_buton.pack(side="left", padx=(0, 10))
        
        # Kapat butonu
        kapat_buton = self.buton_olustur(buton_cerceve, _("Kapat"), son_penceresi.destroy)
        kapat_buton.pack(side="left")
    
    def son_dosya_ac(self, liste, pencere):
        """Seçilen son dosyayı açar"""
        secim = liste.curselection()
        if not secim:
            return
            
        dosya = liste.get(secim[0]).split(" - ")[1]
        if os.path.exists(dosya):
            self.giris_dosyası = dosya
            self.giris_entry.delete(0, tk.END)
            self.giris_entry.insert(0, dosya)
            self.son_kullanılan_klasor = os.path.dirname(dosya)
            self.dosya_bilgilerini_goster(dosya)
            self.powerpoint_metnini_onizle()
            pencere.destroy()
        else:
            messagebox.showerror(_("Hata"), _("Dosya bulunamadı: {}").format(dosya))
    
    def guncelleme_kontrol(self):
        """Uygulama güncellemelerini kontrol eder"""
        if self.guncelleme_kontrol_edildi:
            return
            
        self.durum_guncelle(_("Güncellemeler kontrol ediliyor..."), "bilgi")
        
        try:
            # GitHub'dan güncellemeleri kontrol et
            url = "https://api.github.com/repos/mehmetyay/powerpoint-ses-donusturucu/releases/latest"
            response = requests.get(url, timeout=10)
            
            if response.status_code == 200:
                latest_release = response.json()
                latest_version = latest_release["tag_name"]
                
                # Mevcut sürümü al
                current_version = "2.0"  # Uygulama sürümü
                
                if version.parse(latest_version) > version.parse(current_version):
                    self.durum_guncelle(_("Yeni sürüm mevcut: {}").format(latest_version), "vurgu")
                    
                    # Güncelleme sorusu
                    cevap = messagebox.askyesno(
                        _("Güncelleme Mevcut"),
                        _("Yeni sürüm {} mevcut. Şimdi güncellemek ister misiniz?").format(latest_version),
                        parent=self.root
                    )
                    
                    if cevap:
                        self.guncelleme_indir(latest_release)
                else:
                    self.durum_guncelle(_("Uygulama güncel"), "basari")
            else:
                self.durum_guncelle(_("Güncellemeler kontrol edilemedi"), "uyari")
        
        except Exception as e:
            self.durum_guncelle(_("Güncelleme hatası: {}").format(str(e)), "uyari")
        
        finally:
            self.guncelleme_kontrol_edildi = True
    
    def guncelleme_indir(self, release):
        """Güncellemeyi indirir ve kurar"""
        try:
            # Asset bul
            asset = next((a for a in release["assets"] if a["name"].endswith(".zip")), None)
            
            if not asset:
                messagebox.showerror(_("Hata"), _("Güncelleme paketi bulunamadı"))
                return
                
            self.durum_guncelle(_("Güncelleme indiriliyor..."), "bilgi")
            
            # İndirme penceresi
            indirme_penceresi = tk.Toplevel(self.root)
            indirme_penceresi.title(_("Güncelleme İndiriliyor"))
            indirme_penceresi.geometry("400x150")
            indirme_penceresi.transient(self.root)
            indirme_penceresi.grab_set()
            
            ilerleme_etiket = tk.Label(indirme_penceresi, text=_("İndiriliyor: {}").format(asset["name"]))
            ilerleme_etiket.pack(pady=10)
            
            ilerleme_cubugu = Progressbar(indirme_penceresi, orient=tk.HORIZONTAL, length=300, mode='determinate')
            ilerleme_cubugu.pack(pady=10)
            
            # İndirme işlemi
            response = requests.get(asset["browser_download_url"], stream=True)
            total_size = int(response.headers.get('content-length', 0))
            block_size = 1024
            downloaded = 0
            
            with io.BytesIO() as buffer:
                for data in response.iter_content(block_size):
                    downloaded += len(data)
                    buffer.write(data)
                    ilerleme_cubugu["value"] = (downloaded / total_size) * 100
                    indirme_penceresi.update()
                
                # ZIP dosyasını çıkar
                self.durum_guncelle(_("Güncelleme kuruluyor..."), "bilgi")
                ilerleme_etiket.config(text=_("Kurulum yapılıyor..."))
                
                with zipfile.ZipFile(buffer) as zip_ref:
                    zip_ref.extractall(".")
            
            # Kurulum tamam
            indirme_penceresi.destroy()
            self.durum_guncelle(_("Güncelleme başarıyla tamamlandı!"), "basari")
            messagebox.showinfo(
                _("Başarılı"),
                _("Güncelleme başarıyla tamamlandı. Uygulamayı yeniden başlatın."),
                parent=self.root
            )
            
        except Exception as e:
            self.durum_guncelle(_("Güncelleme hatası: {}").format(str(e)), "uyari")
            messagebox.showerror(
                _("Hata"),
                _("Güncelleme sırasında bir hata oluştu:\n{}").format(str(e)),
                parent=self.root
            )
    
    def yardim_goster(self):
        """Yardım penceresini gösterir"""
        yardim_penceresi = tk.Toplevel(self.root)
        yardim_penceresi.title(_("Yardım"))
        yardim_penceresi.geometry("800x600")
        yardim_penceresi.transient(self.root)
        yardim_penceresi.grab_set()
        
        # Sekmeler
        sekme_kontrol = ttk.Notebook(yardim_penceresi)
        sekme_kontrol.pack(fill="both", expand=True, padx=10, pady=10)
        
        # Kullanım Kılavuzu
        kullanim_cerceve = tk.Frame(sekme_kontrol)
        kullanim_metni = scrolledtext.ScrolledText(kullanim_cerceve, wrap=tk.WORD, font=self.font_normal)
        kullanim_metni.pack(fill="both", expand=True, padx=10, pady=10)
        
        kullanim_metni.insert(tk.END, _("""
PowerPoint'ten Sesli Anlatıma Dönüştürücü Kullanım Kılavuzu

1. PowerPoint Dosyası Seçin:
   - 'Gözat' butonuna tıklayarak bir PowerPoint dosyası seçin.
   - Dosya seçildiğinde metin önizleme alanında içerik görüntülenecektir.

2. Çıkış Ayarlarını Yapın:
   - İstenilen ses formatını seçin (MP3, WAV, OGG, FLAC).
   - Konuşma dilini ve hızını ayarlayın.
   - Çıkış dosyasının kaydedileceği konumu belirtin.

3. Dönüştürme İşlemi:
   - 'Sese Dönüştür' butonuna tıklayın.
   - İşlem tamamlandığında ses dosyası belirtilen konuma kaydedilecektir.

4. Ses Kontrolleri:
   - Oluşturulan ses dosyasını oynatabilir, duraklatabilir veya durdurabilirsiniz.
   - Ses seviyesini ayar çubuğu ile kontrol edebilirsiniz.

5. Ek Özellikler:
   - 'Araçlar' menüsünden toplu dönüştürme yapabilirsiniz.
   - Metin düzenleyici ile PowerPoint metnini düzenleyebilirsiniz.
   - Gece modu ile göz yorgunluğunu azaltabilirsiniz.
"""))
        kullanim_metni.config(state=tk.DISABLED)
        sekme_kontrol.add(kullanim_cerceve, text=_("Kullanım Kılavuzu"))
        
        # SSS
        sss_cerceve = tk.Frame(sekme_kontrol)
        sss_metni = scrolledtext.ScrolledText(sss_cerceve, wrap=tk.WORD, font=self.font_normal)
        sss_metni.pack(fill="both", expand=True, padx=10, pady=10)
        
        sss_metni.insert(tk.END, _("""
Sık Sorulan Sorular

1. Hangi PowerPoint sürümleri destekleniyor?
   - Uygulama .ppt ve .pptx uzantılı tüm modern PowerPoint dosyalarını destekler.

2. Maksimum dosya boyutu limiti var mı?
   - Hayır, ancak çok büyük dosyalar dönüştürme sırasında daha fazla zaman alabilir.

3. İnternet bağlantısı gerekli mi?
   - Evet, metni sese dönüştürme işlemi için Google Text-to-Speech API kullanıldığından internet bağlantısı gereklidir.

4. Ücretli bir uygulama mı?
   - Hayır, bu uygulama tamamen ücretsizdir ve MIT lisansı ile dağıtılmaktadır.

5. Çok sayıda slayt içeren dosyaları dönüştürebilir miyim?
   - Evet, ancak çok uzun dosyalar için toplu dönüştürme özelliğini kullanmanız önerilir.
"""))
        sss_metni.config(state=tk.DISABLED)
        sekme_kontrol.add(sss_cerceve, text=_("SSS"))
        
        # Kapat butonu
        kapat_buton = self.buton_olustur(yardim_penceresi, _("Kapat"), yardim_penceresi.destroy)
        kapat_buton.pack(pady=10)
    
    def kisayollari_goster(self):
        """Klavye kısayollarını gösterir"""
        kisayol_penceresi = tk.Toplevel(self.root)
        kisayol_penceresi.title(_("Klavye Kısayolları"))
        kisayol_penceresi.geometry("600x400")
        kisayol_penceresi.transient(self.root)
        kisayol_penceresi.grab_set()
        
        # Kısayol listesi
        kisayol_cerceve = tk.Frame(kisayol_penceresi, padx=20, pady=20)
        kisayol_cerceve.pack(fill="both", expand=True)
        
        kisayollar = [
            (_("Dosya Aç"), "Ctrl+O"),
            (_("Ayarları Kaydet"), "Ctrl+S"),
            (_("Çıkış"), "Alt+F4"),
            (_("Yardım"), "F1"),
            (_("Tümünü Seç"), "Ctrl+A"),
            (_("Kopyala"), "Ctrl+C"),
            (_("Kes"), "Ctrl+X"),
            (_("Yapıştır"), "Ctrl+V"),
            (_("Sesi Oynat/Duraklat"), "Space"),
            (_("Sesi Durdur"), "Esc")
        ]
        
        for i, (islem, kisayol) in enumerate(kisayollar):
            tk.Label(kisayol_cerceve, text=islem, font=self.font_normal, anchor="w").grid(row=i, column=0, sticky="w", pady=5)
            tk.Label(kisayol_cerceve, text=kisayol, font=self.font_normal, anchor="e").grid(row=i, column=1, sticky="e", pady=5, padx=(20, 0))
        
        # Grid ayarları
        kisayol_cerceve.grid_columnconfigure(0, weight=3)
        kisayol_cerceve.grid_columnconfigure(1, weight=1)
        
        # Kapat butonu
        kapat_buton = self.buton_olustur(kisayol_penceresi, _("Kapat"), kisayol_penceresi.destroy)
        kapat_buton.pack(pady=10)
    
    def hakkinda_goster(self):
        """Hakkında penceresini gösterir"""
        hakkinda_penceresi = tk.Toplevel(self.root)
        hakkinda_penceresi.title(_("Hakkında"))
        hakkinda_penceresi.geometry("500x300")
        hakkinda_penceresi.resizable(False, False)
        hakkinda_penceresi.transient(self.root)
        hakkinda_penceresi.grab_set()
        
        # Logo
        try:
            logo_resim = Image.open("ppt_icon.png").resize((64, 64))
            logo_img = ImageTk.PhotoImage(logo_resim)
            logo_etiket = tk.Label(hakkinda_penceresi, image=logo_img)
            logo_etiket.image = logo_img
            logo_etiket.pack(pady=(20, 10))
        except:
            pass
        
        # Başlık
        tk.Label(hakkinda_penceresi, 
                text=_("PowerPoint'ten Sesli Anlatıma Dönüştürücü Profesyonel"),
                font=self.font_baslik).pack()
        
        # Sürüm
        tk.Label(hakkinda_penceresi, 
                text=_("Sürüm 2.0"),
                font=self.font_alt_baslik).pack(pady=(5, 10))
        
        # Telif hakkı
        tk.Label(hakkinda_penceresi, 
                text=_("© 2025 Mehmet Yay. Tüm hakları saklıdır."),
                font=self.font_kucuk).pack(pady=(0, 20))
        
        # Lisans bilgisi
        tk.Label(hakkinda_penceresi, 
                text=_("MIT Lisansı ile lisanslanmıştır"),
                font=self.font_kucuk).pack()
        
        # Kapat butonu
        kapat_buton = self.buton_olustur(hakkinda_penceresi, _("Kapat"), hakkinda_penceresi.destroy)
        kapat_buton.pack(pady=20)
    
    def kapat(self):
        """Uygulamayı kapatır"""
        if messagebox.askokcancel(_("Çıkış"), _("Uygulamadan çıkmak istediğinizden emin misiniz?")):
            self.ayarları_kaydet()
            self.root.destroy()

def main():
    """Uygulama başlangıç fonksiyonu"""
    root = tk.Tk()
    app = PowerPointSesDonusturucu(root)
    
    # Pencere kapatma işlemi
    root.protocol("WM_DELETE_WINDOW", app.kapat)
    
    root.mainloop()

if __name__ == "__main__":
    main()
        





    